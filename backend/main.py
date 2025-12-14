import os
import uuid
import logging
import random
from datetime import datetime, timedelta
from typing import List, Dict, Optional, Literal

import json
import faiss
from fastapi import (
    FastAPI,
    UploadFile,
    File,
    Form,
    Depends,
    HTTPException,
    Header,
)
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, EmailStr

# File Extractors
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
import zipfile
import csv

# RAG / Embeddings
from sentence_transformers import SentenceTransformer

# Auth / Security
from dotenv import load_dotenv
from passlib.context import CryptContext
from jose import jwt, JWTError

# DB
from motor.motor_asyncio import AsyncIOMotorClient

# Gemini API
from google import genai
from google.oauth2 import id_token
from google.auth.transport import requests as google_requests



# =========================================================
# 🔧 Load Environment
# =========================================================

load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
MONGO_URI = os.getenv("MONGO_URI")
JWT_SECRET = os.getenv("JWT_SECRET", "changeme")
JWT_ALG = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24  # 24 hours

if not GEMINI_API_KEY:
    raise RuntimeError("❌ Missing GEMINI_API_KEY in .env")

if not MONGO_URI:
    raise RuntimeError("❌ Missing MONGO_URI in .env")

genai_client = genai.Client(api_key=GEMINI_API_KEY)



# =========================================================
# ⚙️ FastAPI App Setup
# =========================================================

app = FastAPI(
    title="Smart Campus Assistant",
    description="RAG-powered campus study assistant with multi-format extraction",
    version="3.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*", "http://localhost:3000", "http://127.0.0.1:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("smart-campus")



# =========================================================
# 🗄️ MongoDB Setup
# =========================================================

mongo_client = AsyncIOMotorClient(MONGO_URI)
db = mongo_client["smart_campus"]

users_collection = db["users"]
docs_collection = db["documents"]
shares_collection = db["shares"]



# =========================================================
# 🔐 Security & Auth Helpers
# =========================================================

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
MAX_PASSWORD_LENGTH = 72  # bcrypt limit


def safe_password(p: str) -> str:
    return p[:MAX_PASSWORD_LENGTH]


def hash_password(password: str) -> str:
    return pwd_context.hash(safe_password(password))


def verify_password(plain: str, hashed: str) -> bool:
    return pwd_context.verify(safe_password(plain), hashed)


def create_access_token(user_id: str) -> str:
    payload = {
        "sub": user_id,
        "exp": datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES),
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALG)


async def get_user_by_email(email: str) -> Optional[dict]:
    return await users_collection.find_one({"email": email})


async def get_user_by_id(user_id: str) -> Optional[dict]:
    return await users_collection.find_one({"id": user_id})



# =========================================================
# 🔐 Auth Dependency
# =========================================================

async def get_current_user(Authorization: str = Header(None)):
    """Extract user from Bearer token"""

    if Authorization is None:
        raise HTTPException(401, "Missing Authorization header")

    try:
        scheme, token = Authorization.split()
        if scheme.lower() != "bearer":
            raise ValueError("Invalid scheme")
    except Exception:
        raise HTTPException(401, "Invalid Authorization header format")

    try:
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALG])
        user_id = payload.get("sub")
        if not user_id:
            raise HTTPException(401, "Invalid token payload")
    except JWTError:
        raise HTTPException(401, "Invalid or expired token")

    user = await get_user_by_id(user_id)
    if not user:
        raise HTTPException(401, "User not found")

    return user



# =========================================================
# 🧠 Embeddings + FAISS Setup
# =========================================================

logger.info("Loading sentence-transformer model…")
embedder = SentenceTransformer("all-MiniLM-L6-v2")
EMBED_DIM = embedder.get_sentence_embedding_dimension()
logger.info(f"Embedding model loaded (dim={EMBED_DIM})")

FAISS_INDEXES: Dict[str, faiss.IndexFlatIP] = {}
MAX_CONTEXT_CHARS = 12000



# =========================================================
# 📄 Multi-format Text Extractors
# =========================================================

def extract_text_from_pdf(path: str) -> str:
    try:
        reader = PdfReader(path)
        text = ""
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text += content + "\n"
        return text
    except:
        raise HTTPException(400, "PDF extraction failed")


def extract_text_from_docx(path: str) -> str:
    try:
        doc = DocxDocument(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except:
        raise HTTPException(400, "DOCX extraction failed")


def extract_text_from_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except:
        raise HTTPException(400, "TXT extraction failed")


def extract_text_from_csv(path: str) -> str:
    try:
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            for row in csv.reader(f):
                rows.append(" | ".join(row))
        return "\n".join(rows)
    except:
        raise HTTPException(400, "CSV extraction failed")


def extract_text_from_xlsx(path: str) -> str:
    try:
        wb = load_workbook(path)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text += " | ".join([str(x) for x in row if x]) + "\n"
        return text
    except:
        raise HTTPException(400, "XLS/XLSX extraction failed")


def extract_text_from_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except:
        raise HTTPException(400, "PPTX extraction failed")


def extract_text_from_image(path: str) -> str:
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img)
    except:
        raise HTTPException(400, "Image OCR failed")


def extract_text_from_zip(path: str) -> str:
    try:
        text = ""
        with zipfile.ZipFile(path) as z:
            for filename in z.namelist():
                if filename.endswith((".txt", ".csv")):
                    text += z.read(filename).decode("utf-8", errors="ignore") + "\n"
        return text
    except:
        raise HTTPException(400, "ZIP extraction failed")


def extract_text_from_any(path: str) -> str:
    ext = path.lower().split(".")[-1]

    if ext == "pdf": return extract_text_from_pdf(path)
    if ext == "docx": return extract_text_from_docx(path)
    if ext == "txt": return extract_text_from_txt(path)
    if ext == "csv": return extract_text_from_csv(path)
    if ext in ["xlsx", "xls"]: return extract_text_from_xlsx(path)
    if ext == "pptx": return extract_text_from_pptx(path)
    if ext in ["jpg", "jpeg", "png", "webp"]: return extract_text_from_image(path)
    if ext == "zip": return extract_text_from_zip(path)

    raise HTTPException(400, f"Unsupported file format: {ext}")



# =========================================================
# TEXT CHUNKING + FAISS
# =========================================================

def chunk_text(text: str, max_chars: int = 900) -> List[str]:
    text = text.replace("\n", " ")
    words = text.split()

    chunks, current = [], ""

    for w in words:
        if len(current) + len(w) + 1 <= max_chars:
            current += w + " "
        else:
            if current.strip():
                chunks.append(current.strip())
            current = w + " "

    if current.strip():
        chunks.append(current.strip())

    return chunks


def build_faiss_index_for_doc(doc_id: str, chunks: List[str]):
    embeddings = embedder.encode(chunks, convert_to_numpy=True, normalize_embeddings=True)
    index = faiss.IndexFlatIP(EMBED_DIM)
    index.add(embeddings)
    FAISS_INDEXES[doc_id] = index
    logger.info(f"FAISS index built for doc {doc_id} ({len(chunks)} chunks)")


def retrieve_top_chunks(doc_id: str, query: str, chunks: List[str], top_k=5):
    if doc_id not in FAISS_INDEXES:
        build_faiss_index_for_doc(doc_id, chunks)

    q_emb = embedder.encode([query], convert_to_numpy=True, normalize_embeddings=True)
    index = FAISS_INDEXES[doc_id]

    _, idxs = index.search(q_emb, top_k)
    idxs = idxs[0]

    return [chunks[i] for i in idxs if 0 <= i < len(chunks)]


def trim_context(context: str, max_chars=MAX_CONTEXT_CHARS):
    return context if len(context) <= max_chars else context[-max_chars:]



# =========================================================
# GEMINI API CALL
# =========================================================

def call_gemini(system_prompt: str, user_prompt: str) -> str:
    full = f"{system_prompt}\n\n{user_prompt}"

    try:
        resp = genai_client.models.generate_content(
            model="gemini-2.5-flash",
            contents=full,
        )
        return (resp.text or "").strip()

    except Exception as e:
        logger.exception("Gemini API error: %s", e)
        raise HTTPException(500, "Gemini API error")



# =========================================================
# PYDANTIC MODELS
# =========================================================

class UserBase(BaseModel):
    email: EmailStr
    name: Optional[str] = None


class UserCreate(UserBase):
    password: str


class UserLogin(BaseModel):
    email: EmailStr
    password: str


class UserPublic(UserBase):
    id: str
    provider: str


class TokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    user: UserPublic


class AskRequest(BaseModel):
    doc_id: str
    question: str


class SummarizeRequest(BaseModel):
    doc_id: str
    focus: Optional[str] = None
    mark_type: Optional[Literal["2", "5", "15"]] = "2"


class QuizRequest(BaseModel):
    doc_id: str
    num_questions: int = 5


class GoogleLoginRequest(BaseModel):
    credential: str


class ShareCreateRequest(BaseModel):
    type: Literal["ask", "summary", "quiz"]
    content: dict


class ShareGetRequest(BaseModel):
    code: str



# =========================================================
# AUTH ROUTES
# =========================================================

@app.post("/auth/google", response_model=TokenResponse)
async def google_login(payload: GoogleLoginRequest):
    """
    Google OAuth Login:
    - Verify ID Token
    - Auto-create user if new
    - Return JWT + profile
    """
    try:
        idinfo = id_token.verify_oauth2_token(
            payload.credential,
            google_requests.Request(),
            None,
        )

        email = idinfo.get("email")
        name = idinfo.get("name", email.split("@")[0])

        if not email:
            raise HTTPException(400, "Google token missing email")

        user = await users_collection.find_one({"email": email})

        if not user:
            user = {
                "id": str(uuid.uuid4()),
                "email": email,
                "name": name,
                "password_hash": None,
                "provider": "google",
                "created_at": datetime.utcnow(),
            }
            await users_collection.insert_one(user)

        token = create_access_token(user["id"])

        return TokenResponse(
            access_token=token,
            user=UserPublic(
                id=user["id"],
                email=user["email"],
                name=user.get("name"),
                provider=user.get("provider", "google"),
            ),
        )

    except Exception:
        raise HTTPException(401, "Invalid Google login token")



@app.post("/auth/register", response_model=TokenResponse)
async def register(user_in: UserCreate):
    existing = await get_user_by_email(user_in.email)
    if existing:
        raise HTTPException(400, "Email already registered")

    user_id = str(uuid.uuid4())

    user_doc = {
        "id": user_id,
        "email": user_in.email,
        "name": user_in.name or user_in.email.split("@")[0],
        "password_hash": hash_password(user_in.password),
        "provider": "local",
        "created_at": datetime.utcnow(),
    }

    await users_collection.insert_one(user_doc)

    token = create_access_token(user_id)

    return TokenResponse(
        access_token=token,
        user=UserPublic(
            id=user_id,
            email=user_doc["email"],
            name=user_doc["name"],
            provider="local",
        ),
    )

@app.post("/auth/login", response_model=TokenResponse)
async def login(user_in: UserLogin):
    user = await get_user_by_email(user_in.email)
    if not user:
        raise HTTPException(401, "Invalid credentials")

    hashed = user.get("password_hash")
    if not hashed or not verify_password(user_in.password, hashed):
        raise HTTPException(401, "Invalid credentials")

    token = create_access_token(user["id"])

    return TokenResponse(
        access_token=token,
        user=UserPublic(
            id=user["id"],
            email=user["email"],
            name=user.get("name"),
            provider=user.get("provider", "local"),
        ),
    )


@app.get("/auth/me", response_model=UserPublic)
async def me(current_user: dict = Depends(get_current_user)):
    return UserPublic(
        id=current_user["id"],
        email=current_user["email"],
        name=current_user.get("name"),
        provider=current_user.get("provider", "local"),
    )



# =========================================================
# 🩺 HEALTH CHECK
# =========================================================

@app.get("/")
async def health_check():
    return {
        "status": "ok",
        "message": "Smart Campus Assistant backend running 🚀"
    }


# =========================================================
# 📤 UPLOAD DOCUMENT
# =========================================================

@app.post("/upload")
async def upload_document(
    file: UploadFile = File(...),
    title: str = Form("Untitled"),
    subject: str = Form("General"),
    current_user: dict = Depends(get_current_user),
):
    filename = file.filename
    ext = filename.lower().split(".")[-1]

    allowed = [
        "pdf", "docx", "txt", "csv", "xlsx", "xls",
        "pptx", "jpg", "jpeg", "png", "webp", "zip"
    ]

    if ext not in allowed:
        raise HTTPException(400, f"Unsupported file type: {ext}")

    # Save file with its original extension
    doc_id = str(uuid.uuid4())
    save_path = os.path.join(UPLOAD_DIR, f"{doc_id}.{ext}")

    with open(save_path, "wb") as f:
        f.write(await file.read())

    # Extract text (multi-format)
    raw_text = extract_text_from_any(save_path)

    if not raw_text.strip():
        raise HTTPException(400, "No readable text extracted from file")

    chunks = chunk_text(raw_text)
    created_at = datetime.utcnow()

    doc_entry = {
        "id": doc_id,
        "user_id": current_user["id"],
        "title": title or filename,
        "subject": subject,
        "raw_text": raw_text,
        "chunks": chunks,
        "created_at": created_at,
        "file_ext": ext,
    }

    await docs_collection.insert_one(doc_entry)

    # Build FAISS index
    build_faiss_index_for_doc(doc_id, chunks)

    return {
        "status": "ok",
        "doc_id": doc_id,
        "title": doc_entry["title"],
        "subject": doc_entry["subject"],
        "num_chunks": len(chunks),
        "created_at": created_at.isoformat(),
    }



# =========================================================
# 📚 LIST DOCUMENTS
# =========================================================

@app.get("/documents")
async def list_documents(current_user: dict = Depends(get_current_user)):
    cursor = docs_collection.find(
        {"user_id": current_user["id"]}
    ).sort("created_at", -1)

    docs = []
    async for d in cursor:
        docs.append({
            "id": d["id"],
            "title": d["title"],
            "subject": d["subject"],
            "num_chunks": len(d.get("chunks", [])),
            "created_at": d["created_at"].isoformat(),
        })

    return {"documents": docs}



# =========================================================
# GET DOCUMENT FOR USER
# =========================================================

async def get_doc_for_user(doc_id: str, user_id: str):
    doc = await docs_collection.find_one({"id": doc_id, "user_id": user_id})
    if not doc:
        raise HTTPException(404, "Document not found")
    return doc



# =========================================================
# 🤖 ASK AI (RAG)
# =========================================================

@app.post("/ask")
async def ask_question(
    req: AskRequest,
    current_user: dict = Depends(get_current_user)
):
    doc = await get_doc_for_user(req.doc_id, current_user["id"])
    chunks = doc.get("chunks", [])

    if not chunks:
        raise HTTPException(400, "Document contains no chunks")

    top_chunks = retrieve_top_chunks(req.doc_id, req.question, chunks, top_k=5)
    context = trim_context("\n\n".join(top_chunks))

    system_prompt = (
        "You are a Smart Campus AI tutor.\n"
        "Use ONLY the given CONTEXT from notes.\n"
        "If the answer is missing, reply exactly: 'Not enough information in the notes.'\n"
        "Explain clearly, simply, and exam-focused."
    )

    user_prompt = f"CONTEXT:\n{context}\n\nQUESTION:\n{req.question}"

    answer = call_gemini(system_prompt, user_prompt)

    return {
        "status": "ok",
        "answer": answer,
        "used_chunks": top_chunks,
    }



# =========================================================
# ✍️ SUMMARY GENERATION
# =========================================================

@app.post("/summarize")
async def summarize_document(
    req: SummarizeRequest,
    current_user: dict = Depends(get_current_user),
):
    doc = await get_doc_for_user(req.doc_id, current_user["id"])
    chunks = doc.get("chunks", [])

    if not chunks:
        raise HTTPException(400, "Document contains no chunks")

    selected_chunks = chunks[:15]
    context = trim_context("\n\n".join(selected_chunks))

    # Mark-dependent formatting
    mark = req.mark_type or "2"
    focus_text = req.focus or "important exam topics"

    if mark == "2":
        length_rules = "Write a short 4–5 line answer, crisp and point-wise."
    elif mark == "5":
        length_rules = "Write an 8–10 line answer including definition and key points."
    else:
        length_rules = (
            "Write a 20+ line detailed answer with headings, explanations, "
            "examples, and a conclusion."
        )

    system_prompt = (
        "You are an AI that writes exam-oriented answers using ONLY the student's notes.\n"
        "Never add outside information."
    )

    user_prompt = f"""
CONTEXT:
{context}

TASK:
Provide a {mark}-mark answer on: {focus_text}

RULES:
- {length_rules}
- Use clear exam language
- Use bullets/headings where helpful
- Do NOT say 'summary' or 'AI'
"""

    summary_text = call_gemini(system_prompt, user_prompt)

    return {
        "status": "ok",
        "summary": summary_text,
        "mark_type": mark,
    }



# =========================================================
# 📝 QUIZ GENERATION
# =========================================================

@app.post("/quiz")
async def generate_quiz(
    req: QuizRequest,
    current_user: dict = Depends(get_current_user),
):
    doc = await get_doc_for_user(req.doc_id, current_user["id"])
    chunks = doc.get("chunks", [])

    if not chunks:
        raise HTTPException(400, "Document contains no chunks")

    if not (1 <= req.num_questions <= 20):
        raise HTTPException(400, "num_questions must be between 1 and 20")

    context = trim_context("\n\n".join(chunks[:12]))

    system_prompt = (
        "You are an AI that outputs ONLY VALID JSON.\n"
        "NO markdown. NO commentary."
    )

    user_prompt = f"""
CONTEXT:
{context}

TASK:
Generate exactly {req.num_questions} MCQs in JSON:
[
  {{
    "question": "...",
    "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
    "answer": "A",
    "explanation": "reason"
  }}
]
"""

    raw = call_gemini(system_prompt, user_prompt).strip()

    # Remove accidental markdown code fences
    if raw.startswith("```"):
        raw = raw.replace("```json", "").replace("```", "").strip()

    try:
        quiz_json = json.loads(raw)
    except:
        logger.error("Invalid JSON returned:\n" + raw)
        raise HTTPException(500, "AI returned invalid JSON")

    if not isinstance(quiz_json, list):
        raise HTTPException(500, "Quiz output must be a JSON array")

    return {
        "status": "ok",
        "quiz": quiz_json,
    }
# =========================================================
# 🔗 SHARE CODE GENERATION (6-DIGIT)
# =========================================================

async def generate_unique_code() -> str:
    """
    Generate a unique 6-digit numeric code.
    Retries up to 20 times before failing.
    """
    for _ in range(20):
        code = f"{random.randint(0, 999999):06d}"
        exists = await shares_collection.find_one({"code": code})
        if not exists:
            return code
    raise HTTPException(500, "Failed to generate unique share code")



# =========================================================
# 📤 CREATE SHARE CODE
# =========================================================

@app.post("/share/create")
async def share_create(
    req: ShareCreateRequest,
    current_user: dict = Depends(get_current_user),
):
    """
    Create a 6-digit share code for:
    - Ask AI answer
    - Summary
    - Quiz
    """

    if req.type not in {"ask", "summary", "quiz"}:
        raise HTTPException(400, "Invalid share type")

    code = await generate_unique_code()
    now = datetime.utcnow()

    share_doc = {
        "code": code,
        "type": req.type,
        "content": req.content,
        "owner_user_id": current_user["id"],
        "created_at": now,
    }

    await shares_collection.insert_one(share_doc)

    return {
        "status": "ok",
        "code": code,
        "type": req.type,
        "created_at": now.isoformat(),
    }



# =========================================================
# 📥 RETRIEVE SHARE CODE RESULT
# =========================================================

@app.post("/share/get")
async def share_get(
    req: ShareGetRequest,
    current_user: dict = Depends(get_current_user),
):
    """
    Resolve a 6-digit share code.
    Any logged-in user can access.
    """

    code = req.code.strip()

    if len(code) != 6 or not code.isdigit():
        raise HTTPException(400, "Invalid code format")

    doc = await shares_collection.find_one({"code": code})

    if not doc:
        raise HTTPException(404, "Share code not found")

    return {
        "status": "ok",
        "type": doc["type"],
        "content": doc["content"],
        "code": doc["code"],
    }


